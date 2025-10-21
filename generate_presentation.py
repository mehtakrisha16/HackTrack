# HackTrack Presentation Generator
# This script creates a PowerPoint presentation automatically

# First, install python-pptx library:
# pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)  # 16:9 aspect ratio
prs.slide_height = Inches(5.625)

def add_title_slide(prs, title, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background (you'll need to do this manually or use shapes)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(102, 110, 241)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    footer_frame = footer_box.text_frame
    footer_frame.text = footer
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(16)
    footer_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 110, 241)
    
    # Content
    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(18)

# SLIDE 1: Title
add_title_slide(
    prs,
    "HackTrack",
    "Your Gateway to Tech Opportunities - Connecting Students with Hackathons, Internships & Events Across India",
    "[Your Name] | [Your College] | October 16, 2025\nGitHub: github.com/mehtakrisha16/HackTrack"
)

# SLIDE 2: Problem Statement
add_content_slide(
    prs,
    "The Challenge Students Face",
    [
        "📌 Students miss hackathon & internship opportunities",
        "📌 Information scattered across 50+ platforms",
        "📌 No centralized tracking system",
        "📌 Difficult to find relevant opportunities",
        "📌 Missing application deadlines regularly",
        "",
        "Statistics:",
        "• 70% students miss opportunities due to lack of awareness",
        "• Average student spends 5+ hours searching weekly",
        "• 85% want a centralized platform"
    ]
)

# SLIDE 3: Solution
add_content_slide(
    prs,
    "Our Solution: HackTrack",
    [
        "One Platform. All Opportunities.",
        "",
        "✅ 100+ Real Opportunities",
        "  • 35+ Hackathons (Smart India Hackathon, Google, Microsoft)",
        "  • 40+ Internships (FAANG to Startups)",
        "  • 30+ Events (DevFest, React India, AWS Community)",
        "",
        "✅ Smart Features",
        "  • Advanced filtering (28 options)",
        "  • Real-time countdown timers",
        "  • Direct registration links",
        "",
        '🎯 "Never Miss a Tech Opportunity Again!"'
    ]
)

# SLIDE 4: Key Features
add_content_slide(
    prs,
    "Feature Highlights",
    [
        "🔍 Smart Search & Filters - 28+ options, real-time updates",
        "⏰ Countdown Timers - 5 urgency states, auto-updates",
        "👤 User Profiles - Photo upload, status badges",
        "🔐 Secure Authentication - JWT tokens, Google OAuth",
        "📊 Dashboard Analytics - Personalized recommendations",
        "🔗 Direct Registration - One-click to actual forms"
    ]
)

# SLIDE 5: Tech Stack
add_content_slide(
    prs,
    "Technology Stack",
    [
        "Frontend:",
        "  ⚛️ React 18+ | 🎨 Framer Motion | 🎯 React Router",
        "",
        "Backend:",
        "  🚀 Node.js + Express.js",
        "  🗄️ MongoDB Atlas",
        "  🔐 JWT Authentication",
        "  📤 Multer (File uploads)",
        "  🔒 bcrypt (Security)",
        "",
        "DevOps:",
        "  📦 npm | 🔄 Git & GitHub | 🌐 RESTful APIs"
    ]
)

# SLIDE 6: Homepage Screenshot
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Homepage - First Impression"
content = slide.placeholders[1]
content.text = "📸 INSERT SCREENSHOT HERE: http://localhost:3000/\n\nKey Sections:\n• Hero with gradient\n• Feature cards\n• Statistics\n• Featured opportunities"

# SLIDE 7: Authentication
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Secure Login & Registration"
content = slide.placeholders[1]
content.text = "📸 INSERT SCREENSHOTS: Login & Signup\n\nFeatures:\n• Email/Password auth\n• Google OAuth 2.0\n• JWT tokens (30-day)\n• Password hashing"

# SLIDE 8: Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Personalized Dashboard"
content = slide.placeholders[1]
content.text = "📸 INSERT SCREENSHOT: Dashboard\n\nFeatures:\n• Quick Stats\n• Personalized Recommendations\n• Upcoming Deadlines\n• Real-time Updates"

# SLIDE 9: Hackathons
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = "Hackathons - Find Your Challenge"
content = slide.placeholders[1]
content.text = "📸 INSERT SCREENSHOT: Hackathons page\n\nFeatures:\n• Smart FilterPanel (28 options)\n• Real-time Countdown Timers\n• Event Details\n• Apply Now buttons"

# SLIDE 10: Live Demo
add_content_slide(
    prs,
    "Live Demonstration",
    [
        "🌐 Live URL: http://localhost:3000",
        "📂 GitHub: github.com/mehtakrisha16/HackTrack",
        "",
        "Demo Flow:",
        "1. Homepage → Feature overview",
        "2. Login/Signup → Authentication",
        "3. Dashboard → Personalized view",
        "4. Hackathons → Filters in action",
        "5. Profile → User management",
        "6. Apply Now → Direct registration",
        "",
        '"Let me show you the live application..."'
    ]
)

# SLIDE 11: Thank You / Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
thank_you_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
thank_you_frame = thank_you_box.text_frame
thank_you_frame.text = "Thank You!"
thank_you_para = thank_you_frame.paragraphs[0]
thank_you_para.font.size = Pt(72)
thank_you_para.font.bold = True
thank_you_para.font.color.rgb = RGBColor(102, 110, 241)
thank_you_para.alignment = PP_ALIGN.CENTER

questions_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.5))
questions_frame = questions_box.text_frame
questions_frame.text = "Questions?"
questions_para = questions_frame.paragraphs[0]
questions_para.font.size = Pt(36)
questions_para.alignment = PP_ALIGN.CENTER

contact_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
contact_frame = contact_box.text_frame
contact_frame.text = "📧 [your.email@example.com]\n🔗 GitHub: github.com/mehtakrisha16\n🌐 http://localhost:3000"
contact_para = contact_frame.paragraphs[0]
contact_para.font.size = Pt(16)
contact_para.alignment = PP_ALIGN.CENTER

# Save presentation
prs.save('d:/FINAL/HackTrack_Presentation.pptx')
print("✅ Presentation created successfully!")
print("📁 Saved as: d:/FINAL/HackTrack_Presentation.pptx")
print("\n📝 Next steps:")
print("1. Open the .pptx file")
print("2. Take screenshots of your website")
print("3. Insert screenshots where marked")
print("4. Customize with your details")
print("5. Add more slides from POWERPOINT_SLIDES_CONTENT.txt if needed")
